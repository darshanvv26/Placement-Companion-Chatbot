import os
import json
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from tqdm import tqdm
from transformers import pipeline
from collections import defaultdict
import torch
import pandas as pd
import re
from datetime import datetime

# -------------------------
# Configuration
# -------------------------
BASE_PATH = r"/home/darshan/darshan/darshan/Placement_Data/Placements/Placements"
MMD_DIR = "_mmd_collected"  # OCR results directory

OUTPUT_DIR = "chunks_by_filetype_new"
COMBINED_OUTPUT = "chunks_all.json"
METADATA_OUTPUT = "metadata.json"
SUPPORTED_EXT = {".txt", ".pdf", ".docx", ".pptx"}

# Path to Companies.xlsx - REMOVE trailing slash
COMPANIES_EXCEL_PATH = r"/home/darshan/darshan/darshan/Placement_Data/Placements/Placements/Info/Companies.xlsx"

# Chunking parameters
CHUNK_SIZE_TOKENS = 512
CHUNK_OVERLAP_TOKENS = 64
TOKENS_PER_CHAR = 1 / 4  # Approximate: 1 token ≈ 4 characters
CHUNK_SIZE_CHARS = int(CHUNK_SIZE_TOKENS / TOKENS_PER_CHAR)
CHUNK_OVERLAP_CHARS = int(CHUNK_OVERLAP_TOKENS / TOKENS_PER_CHAR)

# Zero-shot classification model
ZS_MODEL = "MoritzLaurer/DeBERTa-v3-large-mnli-fever-anli-ling-wanli"
LABEL_CONFIDENCE_THRESHOLD = 0.40

# -------------------------
# SECTION DEFINITIONS (Enhanced)
# -------------------------
SECTION_DEFINITIONS = {
    "Eligibility Criteria": {
        "descriptions": [
            "minimum CGPA requirements for applying",
            "academic eligibility and educational qualifications",
            "percentage cutoff criteria for candidates",
            "who is eligible to apply for this position",
            "educational background requirements and restrictions",
            "degree and branch eligibility requirements",
            "no backlogs or academic standing requirements",
            "class 10th and 12th percentage requirements",
            "minimum marks required for eligibility",
            "branch-specific eligibility criteria"
        ],
        "weight": 1.1
    },
    "Job Description": {
        "descriptions": [
            "company overview and background information",
            "about the organization and its mission",
            "role overview and position summary",
            "what the company does and its products",
            "company culture and working environment",
            "team structure and department information",
            "company profile and industry domain",
            "organization vision and values"
        ],
        "weight": 0.85
    },
    "Skill Set Requirements": {
        "descriptions": [
            "required technical skills and programming languages",
            "must-have technical competencies and expertise",
            "preferred skills and technologies knowledge",
            "tools frameworks and libraries required",
            "software proficiency and technical abilities",
            "domain knowledge and specialization required",
            "programming languages and coding skills",
            "technical certifications and qualifications"
        ],
        "weight": 1.0
    },
    "Key Responsibilities": {
        "descriptions": [
            "day-to-day job responsibilities and duties",
            "what you will be doing in this role",
            "core job functions and deliverables",
            "primary tasks and work activities",
            "responsibilities you will handle daily",
            "job duties and work expectations",
            "key performance indicators and outcomes",
            "project responsibilities and deliverables"
        ],
        "weight": 0.85
    },
    "Selection Rounds": {
        "descriptions": [
            "interview process and hiring rounds",
            "online test and assessment details",
            "technical rounds and coding tests",
            "HR interview and managerial rounds",
            "selection procedure and evaluation stages",
            "test platform and examination format",
            "number of questions and test duration",
            "shortlisting process and next rounds",
            "aptitude test and reasoning assessment",
            "group discussion and presentation rounds"
        ],
        "weight": 1.05
    },
    "Compensation": {
        "descriptions": [
            "salary package and CTC details",
            "stipend amount for internship",
            "compensation structure and pay scale",
            "remuneration and financial benefits",
            "salary range and package offered",
            "monthly stipend or annual salary",
            "variable pay and bonus structure",
            "cost to company breakdown"
        ],
        "weight": 1.0
    },
    "Benefits": {
        "descriptions": [
            "what the company offers to employees",
            "perks and benefits package details",
            "employee advantages and rewards",
            "additional benefits beyond salary",
            "work environment and facilities",
            "health insurance and medical benefits",
            "learning and development opportunities",
            "work-life balance initiatives"
        ],
        "weight": 0.75
    },
    "Location": {
        "descriptions": [
            "work location and office address",
            "where the job is located geographically",
            "office location or remote work details",
            "workplace location information",
            "city and state of work location",
            "remote or hybrid work options",
            "multiple office locations available"
        ],
        "weight": 0.7
    },
    "Program Details": {
        "descriptions": [
            "internship program structure and duration",
            "training program details and timeline",
            "program overview and learning opportunities",
            "internship duration and start date",
            "program schedule and milestones",
            "full-time conversion opportunities",
            "mentorship and guidance provided"
        ],
        "weight": 0.8
    },
    "Drive Details": {
        "descriptions": [
            "placement drive schedule and dates",
            "drive timeline and process flow",
            "problem statement and assignment details",
            "drive description and important dates",
            "registration deadline and drive date",
            "drive venue and timing information",
            "pre-placement talk details"
        ],
        "weight": 0.9
    },
    "Company Statistics": {
        "descriptions": [
            "number of positions available",
            "hiring statistics and openings",
            "company size and employee count",
            "previous year placement statistics",
            "recruitment numbers and targets"
        ],
        "weight": 0.8
    },
    "MMD": {
        "descriptions": [
            "multimodal document content",
            "OCR extracted text from images",
            "visual document information",
            "scanned document content"
        ],
        "weight": 0.9
    },
    "Other": {
        "descriptions": ["miscellaneous information", "general content", "uncategorized"],
        "weight": 0.1
    }
}

CANDIDATE_LABELS = list(SECTION_DEFINITIONS.keys())

# -------------------------
# Metadata Extraction Functions
# -------------------------
def extract_compensation_from_text(text):
    """Extract both stipend and CTC information from text."""
    compensation = {
        "stipend": None,
        "ctc": None,
        "salary_type": None  # "monthly_stipend", "annual_ctc", "hourly"
    }
    
    # Stipend patterns (monthly/hourly)
    stipend_patterns = [
        r'stipend.*?₹\s*(\d+(?:,\d+)*(?:\.\d+)?)\s*(?:k|K|thousand)?(?:\s*per\s*month)?',
        r'monthly\s+stipend.*?₹\s*(\d+(?:,\d+)*(?:\.\d+)?)',
        r'₹\s*(\d+(?:,\d+)*(?:\.\d+)?)\s*(?:k|K)?\s*per\s*month',
        r'₹\s*(\d+(?:,\d+)*(?:\.\d+)?)\s*(?:k|K)?\s*monthly',
    ]
    
    # CTC patterns (annual)
    ctc_patterns = [
        r'(?:ctc|package|annual\s+package).*?₹\s*(\d+(?:,\d+)*(?:\.\d+)?)\s*(?:lpa|LPA|lakhs?\s*per\s*annum)',
        r'₹\s*(\d+(?:,\d+)*(?:\.\d+)?)\s*(?:lpa|LPA)',
        r'(?:ctc|package).*?₹\s*(\d+(?:,\d+)*(?:\.\d+)?)\s*(?:L|lakhs?)',
        r'(\d+(?:,\d+)*(?:\.\d+)?)\s*(?:lpa|LPA)',
    ]
    
    # Salary patterns (generic)
    salary_patterns = [
        r'salary.*?₹\s*(\d+(?:,\d+)*(?:\.\d+)?)\s*(?:k|K|thousand|lpa|LPA)?',
        r'remuneration.*?₹\s*(\d+(?:,\d+)*(?:\.\d+)?)',
    ]
    
    # Try to extract stipend
    for pattern in stipend_patterns:
        matches = re.findall(pattern, text, re.IGNORECASE)
        if matches:
            try:
                value = float(matches[0].replace(',', ''))
                compensation["stipend"] = value
                compensation["salary_type"] = "monthly_stipend"
                break
            except:
                continue
    
    # Try to extract CTC
    for pattern in ctc_patterns:
        matches = re.findall(pattern, text, re.IGNORECASE)
        if matches:
            try:
                value = float(matches[0].replace(',', ''))
                compensation["ctc"] = value
                compensation["salary_type"] = "annual_ctc"
                break
            except:
                continue
    
    # Fallback to generic salary if neither found
    if not compensation["stipend"] and not compensation["ctc"]:
        for pattern in salary_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE)
            if matches:
                try:
                    value = float(matches[0].replace(',', ''))
                    # Try to determine if it's stipend or CTC based on context
                    if re.search(r'per\s+month|monthly', text, re.IGNORECASE):
                        compensation["stipend"] = value
                        compensation["salary_type"] = "monthly_stipend"
                    elif re.search(r'per\s+annum|annual|lpa', text, re.IGNORECASE):
                        compensation["ctc"] = value
                        compensation["salary_type"] = "annual_ctc"
                    else:
                        # Default to stipend for smaller values, CTC for larger
                        if value < 100:  # Less than 100 likely monthly
                            compensation["stipend"] = value
                            compensation["salary_type"] = "monthly_stipend"
                        else:  # 100+ likely annual
                            compensation["ctc"] = value
                            compensation["salary_type"] = "annual_ctc"
                    break
                except:
                    continue
    
    return compensation

def extract_stipend_from_text(text):
    """Extract stipend/salary information from text (kept for backward compatibility)."""
    comp = extract_compensation_from_text(text)
    return comp["stipend"] or comp["ctc"]

def extract_cgpa_from_text(text):
    """Extract CGPA requirements from text."""
    patterns = [
        r'(?:minimum|min\.?)\s*(?:cgpa|CGPA).*?(\d+(?:\.\d+)?)',
        r'(?:cgpa|CGPA).*?(\d+(?:\.\d+)?)',
        r'(\d+(?:\.\d+)?)\s*(?:cgpa|CGPA)',
    ]
    
    for pattern in patterns:
        matches = re.findall(pattern, text, re.IGNORECASE)
        if matches:
            try:
                cgpa = float(matches[0])
                if 0 <= cgpa <= 10:
                    return cgpa
            except:
                continue
    return None

def extract_year_from_folder(company_name):
    """Extract year from company folder name."""
    match = re.search(r'(\d{4})', company_name)
    return int(match.group(1)) if match else None

def load_companies_excel():
    """Load and parse Companies.xlsx for year-wise company visits."""
    try:
        # Remove trailing slash from path
        excel_path = COMPANIES_EXCEL_PATH.rstrip('/')
        
        if not os.path.exists(excel_path):
            print(f"⚠️ Companies.xlsx not found at: {excel_path}")
            return {}
        
        # Read Excel without header
        df = pd.read_excel(excel_path, header=None)
        
        print(f"📊 Excel loaded: {len(df)} rows, {len(df.columns)} columns")
        
        # Parse the structure
        company_visits = {}
        
        # First row contains years (2023, 2024, 2025) and other headers (BDA, AIML)
        first_row = df.iloc[0]
        year_columns = {}
        
        # Identify year columns
        for col_idx, val in enumerate(first_row):
            if pd.notna(val):
                val_str = str(val).strip()
                # Check if it's a 4-digit year
                if val_str.isdigit() and len(val_str) == 4:
                    year = int(val_str)
                    if 2020 <= year <= 2030:  # Reasonable year range
                        year_columns[year] = col_idx
                        print(f"   Found year column: {year} at index {col_idx}")
        
        # Find BDA and AIML column indices (for 2025 eligibility data)
        bda_col = None
        aiml_col = None
        placed_col_offset = 1  # "No of Placed" is typically 1 column after
        
        for col_idx, val in enumerate(first_row):
            if pd.notna(val):
                val_str = str(val).strip().upper()
                if val_str == 'BDA':
                    bda_col = col_idx
                    print(f"   Found BDA column at index {col_idx}")
                elif val_str == 'AIML':
                    aiml_col = col_idx
                    print(f"   Found AIML column at index {col_idx}")
        
        if not year_columns:
            print("⚠️ No year columns found in Excel!")
            return {}
        
        # Process each row starting from row 1
        for row_idx in range(1, len(df)):
            row = df.iloc[row_idx]
            
            # Check each year column for company names
            for year, col_idx in year_columns.items():
                if col_idx >= len(row):
                    continue
                
                company_name = row.iloc[col_idx]
                
                # Skip empty cells
                if pd.isna(company_name):
                    continue
                
                company_name = str(company_name).strip()
                
                # Skip if empty string or looks like a number/date
                if not company_name or company_name.replace('.', '').isdigit():
                    continue
                
                # Initialize company entry if not exists
                if company_name not in company_visits:
                    company_visits[company_name] = {
                        'years_visited': [],
                        'visit_count': 0,
                        'is_recurring': False,
                        'bda_eligible': {},
                        'aiml_eligible': {},
                        'bda_placed': {},
                        'aiml_placed': {}
                    }
                
                # Add year to visits (avoid duplicates)
                if year not in company_visits[company_name]['years_visited']:
                    company_visits[company_name]['years_visited'].append(year)
                    company_visits[company_name]['visit_count'] += 1
            
            # Extract BDA/AIML eligibility for 2025 (if row has data in those columns)
            if bda_col and bda_col < len(row):
                bda_value = row.iloc[bda_col]
                if pd.notna(bda_value):
                    bda_status = str(bda_value).strip()
                    
                    # Try to find company name in 2025 column for this row
                    year_2025_col = year_columns.get(2025)
                    if year_2025_col and year_2025_col < len(row):
                        company_name = row.iloc[year_2025_col]
                        if pd.notna(company_name):
                            company_name = str(company_name).strip()
                            if company_name and company_name in company_visits:
                                company_visits[company_name]['bda_eligible'][2025] = bda_status
                                
                                # Get placement count
                                placed_col = bda_col + placed_col_offset
                                if placed_col < len(row):
                                    placed_count = row.iloc[placed_col]
                                    if pd.notna(placed_count):
                                        try:
                                            company_visits[company_name]['bda_placed'][2025] = int(placed_count)
                                        except:
                                            pass
            
            # Extract AIML eligibility
            if aiml_col and aiml_col < len(row):
                aiml_value = row.iloc[aiml_col]
                if pd.notna(aiml_value):
                    aiml_status = str(aiml_value).strip()
                    
                    # Try to find company name in 2025 column for this row
                    year_2025_col = year_columns.get(2025)
                    if year_2025_col and year_2025_col < len(row):
                        company_name = row.iloc[year_2025_col]
                        if pd.notna(company_name):
                            company_name = str(company_name).strip()
                            if company_name and company_name in company_visits:
                                company_visits[company_name]['aiml_eligible'][2025] = aiml_status
                                
                                # Get placement count
                                placed_col = aiml_col + placed_col_offset
                                if placed_col < len(row):
                                    placed_count = row.iloc[placed_col]
                                    if pd.notna(placed_count):
                                        try:
                                            company_visits[company_name]['aiml_placed'][2025] = int(placed_count)
                                        except:
                                            pass
        
        # Mark recurring companies (visited 2+ years)
        recurring_count = 0
        for company in company_visits:
            # Sort years visited
            company_visits[company]['years_visited'].sort()
            
            # Mark as recurring if visited in 2 or more years
            if company_visits[company]['visit_count'] >= 2:
                company_visits[company]['is_recurring'] = True
                recurring_count += 1
        
        print(f"✅ Parsed {len(company_visits)} unique companies from Excel")
        print(f"   Recurring companies: {recurring_count}")
        print(f"   Years tracked: {sorted(year_columns.keys())}")
        
        # Debug: Print first few companies
        if company_visits:
            print(f"\n   Sample companies:")
            for i, (company, info) in enumerate(list(company_visits.items())[:3]):
                print(f"     - {company}: Years {info['years_visited']}, Recurring: {info['is_recurring']}")
        
        return company_visits
        
    except Exception as e:
        print(f"⚠️ Error parsing Companies.xlsx: {e}")
        import traceback
        traceback.print_exc()
        return {}

# -------------------------
# File Readers
# -------------------------
def read_txt(file_path):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def read_pdf(file_path):
    text = ""
    try:
        with fitz.open(file_path) as doc:
            for page in doc:
                text += page.get_text("text")
    except Exception as e:
        print(f"⚠️ PDF read error: {e}")
    return text

def read_docx(file_path):
    try:
        doc = Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    except Exception as e:
        print(f"⚠️ DOCX read error: {e}")
        return ""

def read_pptx(file_path):
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"⚠️ PPTX read error: {e}")
    return text

def read_mmd(file_path):
    """Read MMD/OCR result files."""
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:
        print(f"⚠️ MMD read error: {e}")
        return ""

def extract_text_from_file(file_path):
    ext = os.path.splitext(file_path)[-1].lower()
    readers = {
        ".txt": read_txt, 
        ".pdf": read_pdf, 
        ".docx": read_docx, 
        ".pptx": read_pptx,
        ".mmd": read_mmd
    }
    return readers.get(ext, lambda _: "")(file_path)

# -------------------------
# Chunking Logic
# -------------------------
def char_chunks(text):
    start = 0
    while start < len(text):
        end = start + CHUNK_SIZE_CHARS
        chunk = text[start:end].strip()
        if chunk:
            yield chunk
        if end >= len(text):
            break
        start = end - CHUNK_OVERLAP_CHARS

# -------------------------
# Classification
# -------------------------
def classify_section(classifier, text, fallback="Other"):
    try:
        result = classifier(
            text, CANDIDATE_LABELS,
            hypothesis_template="This text is about {}.", multi_label=False
        )
        if not result or "labels" not in result:
            return fallback, 0.0
        weighted = {l: s * SECTION_DEFINITIONS[l]["weight"] for l, s in zip(result["labels"], result["scores"])}
        best = max(weighted, key=weighted.get)
        score = result["scores"][result["labels"].index(best)]
        return (best, score) if score >= LABEL_CONFIDENCE_THRESHOLD else (fallback, score)
    except Exception as e:
        print(f"⚠️ Classification error: {e}")
        return fallback, 0.0

# -------------------------
# MMD/OCR File Collection
# -------------------------
def collect_mmd_files():
    """Collect all MMD/OCR result files from _mmd_collected directory."""
    mmd_files = []
    if not os.path.exists(MMD_DIR):
        print(f"⚠️ MMD directory not found: {MMD_DIR}")
        return mmd_files
    
    for root, _, files in os.walk(MMD_DIR):
        for name in files:
            if name.lower() == "result.mmd":
                mmd_files.append(os.path.join(root, name))
    
    return mmd_files

def parse_mmd_metadata(file_path, mmd_base=MMD_DIR):
    """Extract metadata from MMD file path."""
    rel = os.path.relpath(file_path, mmd_base)
    parts = rel.split(os.sep)
    company = parts[0] if parts else "Unknown"
    role = parts[1] if len(parts) > 1 else "General"
    filename = os.path.basename(file_path)
    return {"company": company, "role": role, "filename": filename, "file_type": ".mmd"}

# -------------------------
# File Scanning & Metadata
# -------------------------
def scan_all_companies(base_path=BASE_PATH):
    """Scan all company directories and return file paths."""
    all_files = []
    skip_folders = {"Info"}

    # Auto-detect wrapper folder like "Placements"
    subdirs = [d for d in os.listdir(base_path) if os.path.isdir(os.path.join(base_path, d))]
    if len(subdirs) == 1 and subdirs[0].lower() in {"placements", "companies", "data"}:
        base_path = os.path.join(base_path, subdirs[0])
        print(f"📁 Auto-detected wrapper folder, switching BASE_PATH → {base_path}")

    for company_folder in os.listdir(base_path):
        company_path = os.path.join(base_path, company_folder)
        if not os.path.isdir(company_path) or company_folder in skip_folders:
            continue
        for root, _, files in os.walk(company_path):
            for file in files:
                ext = os.path.splitext(file)[-1].lower()
                if ext in SUPPORTED_EXT:
                    all_files.append(os.path.join(root, file))
    return all_files

def parse_file_metadata(file_path, base_path=BASE_PATH):
    rel = os.path.relpath(file_path, base_path)
    parts = rel.split(os.sep)
    company = parts[0] if parts else "Unknown"
    role = parts[1] if len(parts) > 2 else "General"
    filename = os.path.basename(file_path)
    ext = os.path.splitext(filename)[-1].lower()
    return {"company": company, "role": role, "filename": filename, "file_type": ext}

def build_global_metadata(company_chunks, companies_info):
    """Build comprehensive metadata for the chatbot."""
    metadata = {
        "total_companies": len(company_chunks),
        "companies": sorted(list(company_chunks.keys())),
        "generation_timestamp": datetime.now().isoformat(),
        "statistics": {
            "total_chunks": 0,
            "sections_distribution": defaultdict(int),
            "file_types_distribution": defaultdict(int)
        },
        "company_details": {},
        "compensation": {
            "stipends": [],
            "ctc_packages": [],
            "avg_stipend": None,
            "highest_stipend": None,
            "lowest_stipend": None,
            "avg_ctc": None,
            "highest_ctc": None,
            "lowest_ctc": None
        },
        "eligibility": {
            "cgpa_requirements": [],
            "avg_cgpa": None
        },
        "years": defaultdict(list),
        "company_visits": companies_info,
        "recurring_companies": [c for c, info in companies_info.items() if info.get('is_recurring', False)]
    }
    
    all_stipends = []
    all_ctcs = []
    all_cgpas = []
    
    for company, chunks in company_chunks.items():
        company_meta = {
            "total_chunks": len(chunks),
            "roles": set(),
            "file_types": set(),
            "sections": defaultdict(int),
            "compensation": {
                "stipend": None,
                "ctc": None,
                "salary_type": None
            },
            "cgpa_requirement": None,
            "year": extract_year_from_folder(company)
        }
        
        # Match with Excel data
        excel_match = None
        for excel_company, info in companies_info.items():
            if (excel_company.lower() in company.lower() or 
                company.lower() in excel_company.lower() or
                excel_company.lower().replace(' ', '') in company.lower().replace(' ', '')):
                excel_match = excel_company
                break
        
        if excel_match:
            company_meta["excel_info"] = {
                "exact_name": excel_match,
                "years_visited": companies_info[excel_match].get('years_visited', []),
                "visit_count": companies_info[excel_match].get('visit_count', 0),
                "is_recurring": companies_info[excel_match].get('is_recurring', False),
                "bda_eligible": companies_info[excel_match].get('bda_eligible', {}),
                "aiml_eligible": companies_info[excel_match].get('aiml_eligible', {}),
                "bda_placed": companies_info[excel_match].get('bda_placed', {}),
                "aiml_placed": companies_info[excel_match].get('aiml_placed', {})
            }
        
        full_text = ""
        
        for chunk in chunks:
            metadata["statistics"]["total_chunks"] += 1
            metadata["statistics"]["sections_distribution"][chunk["section"]] += 1
            metadata["statistics"]["file_types_distribution"][chunk["file_type"]] += 1
            
            company_meta["roles"].add(chunk["role"])
            company_meta["file_types"].add(chunk["file_type"])
            company_meta["sections"][chunk["section"]] += 1
            
            full_text += " " + chunk["content"]
        
        # Extract compensation (both stipend and CTC)
        compensation = extract_compensation_from_text(full_text)
        company_meta["compensation"] = compensation
        
        if compensation["stipend"]:
            all_stipends.append(compensation["stipend"])
        if compensation["ctc"]:
            all_ctcs.append(compensation["ctc"])
        
        # Extract CGPA
        cgpa = extract_cgpa_from_text(full_text)
        if cgpa:
            company_meta["cgpa_requirement"] = cgpa
            all_cgpas.append(cgpa)
        
        # Convert sets to lists for JSON serialization
        company_meta["roles"] = list(company_meta["roles"])
        company_meta["file_types"] = list(company_meta["file_types"])
        company_meta["sections"] = dict(company_meta["sections"])
        
        metadata["company_details"][company] = company_meta
        
        # Group by year
        if company_meta["year"]:
            metadata["years"][company_meta["year"]].append(company)
    
    # Calculate stipend statistics
    if all_stipends:
        metadata["compensation"]["stipends"] = all_stipends
        metadata["compensation"]["avg_stipend"] = f"₹{sum(all_stipends)/len(all_stipends):.1f} K per month"
        metadata["compensation"]["highest_stipend"] = f"₹{max(all_stipends):.1f} K per month"
        metadata["compensation"]["lowest_stipend"] = f"₹{min(all_stipends):.1f} K per month"
    
    # Calculate CTC statistics
    if all_ctcs:
        metadata["compensation"]["ctc_packages"] = all_ctcs
        metadata["compensation"]["avg_ctc"] = f"₹{sum(all_ctcs)/len(all_ctcs):.1f} LPA"
        metadata["compensation"]["highest_ctc"] = f"₹{max(all_ctcs):.1f} LPA"
        metadata["compensation"]["lowest_ctc"] = f"₹{min(all_ctcs):.1f} LPA"
    
    # Calculate CGPA statistics
    if all_cgpas:
        metadata["eligibility"]["cgpa_requirements"] = all_cgpas
        metadata["eligibility"]["avg_cgpa"] = round(sum(all_cgpas)/len(all_cgpas), 2)
        metadata["eligibility"]["min_cgpa"] = min(all_cgpas)
        metadata["eligibility"]["max_cgpa"] = max(all_cgpas)
    
    # Convert defaultdicts to regular dicts
    metadata["statistics"]["sections_distribution"] = dict(metadata["statistics"]["sections_distribution"])
    metadata["statistics"]["file_types_distribution"] = dict(metadata["statistics"]["file_types_distribution"])
    metadata["years"] = dict(metadata["years"])
    
    return metadata

# -------------------------
# Main Processing
# -------------------------
def create_chunks_by_company():
    print("=" * 60)
    print("PLACEMENT COMPANION - COMPANY-CENTRIC CHUNKING")
    print("=" * 60)

    print("\n📁 Scanning placement folders...")
    all_files = scan_all_companies(BASE_PATH)
    print(f"✅ Found {len(all_files)} regular files")

    # Collect MMD/OCR files
    print("\n🔍 Scanning MMD/OCR results...")
    mmd_files = collect_mmd_files()
    print(f"✅ Found {len(mmd_files)} MMD files")

    print("\n📊 Loading Companies.xlsx metadata...")
    companies_info = load_companies_excel()
    print(f"✅ Loaded info for {len(companies_info)} companies from Excel")

    print("\n🧠 Loading DeBERTa-v3-large classifier...")
    device = 0 if torch.cuda.is_available() else -1
    classifier = pipeline("zero-shot-classification", model=ZS_MODEL, device=device)
    print(f"✅ Classifier ready on {'GPU' if device == 0 else 'CPU'}")

    print("\n📝 Generating chunks from regular files...")
    company_chunks = defaultdict(list)
    file_types_summary = defaultdict(lambda: {"files": 0, "chunks": 0})
    total_files = 0
    total_chunks = 0

    for f_path in tqdm(all_files, desc="Processing Regular Files"):
        meta = parse_file_metadata(f_path)
        company, role, filename, ext = meta["company"], meta["role"], meta["filename"], meta["file_type"]
        text = extract_text_from_file(f_path)
        if not text.strip():
            continue

        total_files += 1
        file_types_summary[ext]["files"] += 1

        for idx, chunk_text in enumerate(char_chunks(text), start=1):
            section, conf = classify_section(classifier, chunk_text)
            chunk = {
                "chunk_id": idx,
                "section": section,
                "confidence": round(conf, 3),
                "content": chunk_text,
                "company": company,
                "role": role,
                "filename": filename,
                "file_type": ext
            }
            company_chunks[company].append(chunk)
            total_chunks += 1
            file_types_summary[ext]["chunks"] += 1

    # Process MMD files
    print("\n📝 Generating chunks from MMD/OCR files...")
    mmd_added = 0
    for mmd_path in tqdm(mmd_files, desc="Processing MMD Files"):
        meta = parse_mmd_metadata(mmd_path)
        company, role, filename = meta["company"], meta["role"], meta["filename"]
        
        text = read_mmd(mmd_path)
        if not text.strip():
            continue

        total_files += 1
        file_types_summary[".mmd"]["files"] += 1

        # For MMD files, create chunks but mark as MMD section
        for chunk_text in char_chunks(text):
            # Optionally classify MMD content too, or just mark as MMD
            section, conf = classify_section(classifier, chunk_text, fallback="MMD")
            chunk = {
                "chunk_id": len(company_chunks[company]) + 1,
                "section": section if section != "Other" else "MMD",
                "confidence": round(conf, 3),
                "content": chunk_text,
                "company": company,
                "role": role,
                "filename": filename,
                "file_type": ".mmd"
            }
            company_chunks[company].append(chunk)
            total_chunks += 1
            mmd_added += 1
            file_types_summary[".mmd"]["chunks"] += 1

    print(f"✅ Added {mmd_added} chunks from MMD files")

    # Convert to company-keyed structure
    chunks_structured = {company: chunks for company, chunks in company_chunks.items()}

    print("\n📊 Building comprehensive metadata...")
    metadata = build_global_metadata(company_chunks, companies_info)

    print("\n💾 Saving combined JSON...")
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    combined_data = {
        "file_type": "all",
        "total_files": total_files,
        "total_chunks": total_chunks,
        "file_types": dict(file_types_summary),
        "metadata": metadata,
        "chunks": chunks_structured
    }

    out_path = os.path.join(OUTPUT_DIR, COMBINED_OUTPUT)
    with open(out_path, "w", encoding="utf-8") as f:
        json.dump(combined_data, f, indent=2, ensure_ascii=False)

    print(f"✅ Saved: {out_path}")
    
    # Save standalone metadata file
    metadata_path = os.path.join(OUTPUT_DIR, METADATA_OUTPUT)
    with open(metadata_path, "w", encoding="utf-8") as f:
        json.dump(metadata, f, indent=2, ensure_ascii=False)
    
    print(f"✅ Saved metadata: {metadata_path}")
    print(f"Companies processed: {len(company_chunks)} | Total chunks: {total_chunks}")
    return combined_data, metadata

# -------------------------
# Main Execution
# -------------------------
if __name__ == "__main__":
    data, metadata = create_chunks_by_company()
    print("\n" + "=" * 60)
    print("CHUNKING SUMMARY")
    print("=" * 60)
    print(f"Total Files: {data['total_files']}")
    print(f"Total Chunks: {data['total_chunks']}")
    print(f"Total Companies: {metadata['total_companies']}")
    print(f"Sample Companies: {metadata['companies'][:5]}")
    
    # Display both stipend and CTC stats
    if metadata['compensation']['avg_stipend']:
        print(f"\n💰 Stipend Statistics:")
        print(f"   Average: {metadata['compensation']['avg_stipend']}")
        print(f"   Highest: {metadata['compensation']['highest_stipend']}")
        print(f"   Lowest: {metadata['compensation']['lowest_stipend']}")
    
    if metadata['compensation']['avg_ctc']:
        print(f"\n💼 CTC Statistics:")
        print(f"   Average: {metadata['compensation']['avg_ctc']}")
        print(f"   Highest: {metadata['compensation']['highest_ctc']}")
        print(f"   Lowest: {metadata['compensation']['lowest_ctc']}")
    
    if metadata['eligibility'].get('avg_cgpa'):
        print(f"\n📚 CGPA Requirements:")
        print(f"   Average: {metadata['eligibility']['avg_cgpa']}")
        print(f"   Range: {metadata['eligibility']['min_cgpa']} - {metadata['eligibility']['max_cgpa']}")
    
    print(f"\n📅 Years Coverage: {list(metadata['years'].keys())}")
    print(f"📂 File Types: {list(data['file_types'].keys())}")
    
    if metadata.get('recurring_companies'):
        print(f"\n🔄 Recurring Companies: {len(metadata['recurring_companies'])}")
        print(f"   Examples: {metadata['recurring_companies'][:5]}")
